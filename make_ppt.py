"""
「책잇다」 SSAFY 관통 PJT 발표용 PPT 생성 스크립트
(STATUS.md 기반 최신화 버전)
실행: python make_ppt.py
출력: 책잇다_발표.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── 컬러 팔레트 ──────────────────────────────────────────────────────────────
C_DARK    = RGBColor(0x1A, 0x1A, 0x2E)
C_MID     = RGBColor(0x16, 0x21, 0x3E)
C_ACCENT  = RGBColor(0x0F, 0x3C, 0x78)
C_POINT   = RGBColor(0x45, 0xA2, 0x9D)
C_GOLD    = RGBColor(0xE8, 0xB8, 0x60)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT   = RGBColor(0xE8, 0xF0, 0xFE)
C_GRAY    = RGBColor(0xB0, 0xBE, 0xC5)
C_RED     = RGBColor(0xE5, 0x53, 0x4B)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

def add_rect(slide, l, t, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    return shape

def add_text(slide, text, l, t, w, h, size=18, bold=False, color=C_WHITE, align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb

def bg_dark(slide, color=C_DARK):
    add_rect(slide, 0, 0, W, H, fill=color)

def accent_bar(slide, color=C_POINT, thickness=Inches(0.06)):
    add_rect(slide, 0, 0, W, thickness, fill=color)

def slide_cover():
    s = prs.slides.add_slide(BLANK)
    bg_dark(s)
    add_rect(s, Inches(8.5), 0, Inches(4.83), H, fill=C_ACCENT)
    add_rect(s, Inches(8.3), Inches(0.5), Inches(0.06), Inches(6.5), fill=C_GOLD)
    accent_bar(s, C_POINT, Inches(0.08))

    badge = add_rect(s, Inches(0.6), Inches(1.2), Inches(2.1), Inches(0.42), fill=C_POINT)
    badge.line.fill.background()
    add_text(s, "📚  책잇다", Inches(0.6), Inches(1.18), Inches(2.1), Inches(0.45), size=13, bold=True, align=PP_ALIGN.CENTER)

    add_text(s, "책잇다", Inches(0.6), Inches(1.7), Inches(7.5), Inches(1.5), size=72, bold=True)

    m1 = add_rect(s, Inches(0.6), Inches(3.25), Inches(3.6), Inches(0.45), fill=C_ACCENT, line=C_GOLD)
    add_text(s, "책 · 도서관 · 사람을 잇다", Inches(0.65), Inches(3.28), Inches(3.5), Inches(0.42), size=13, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)
    
    m2 = add_rect(s, Inches(4.4), Inches(3.25), Inches(2.8), Inches(0.45), fill=C_ACCENT, line=C_POINT)
    add_text(s, "도서관에 책이 있다", Inches(4.45), Inches(3.28), Inches(2.7), Inches(0.42), size=13, bold=True, color=C_POINT, align=PP_ALIGN.CENTER)

    add_text(s, "지금 내 도서관에서 빌릴 수 있는\n취향 맞춤 도서 추천 서비스", Inches(0.6), Inches(3.85), Inches(7.5), Inches(1.3), size=20, color=C_LIGHT)
    add_text(s, "완독 직후 • 도서관 방문 전  ─  딱 그 순간", Inches(0.6), Inches(5.2), Inches(7.5), Inches(0.55), size=14, color=C_POINT, italic=True)

    add_text(s, "SSAFY 관통 PJT · 2인 팀", Inches(8.7), Inches(1.5), Inches(4.3), Inches(0.6), size=14, color=C_GOLD, bold=True)
    add_text(s, "Django × Vue\n하이브리드 AI 추천 엔진\n도서관 실시간 가용성 매핑", Inches(8.7), Inches(2.2), Inches(4.3), Inches(2.0), size=17)
    add_text(s, "정보나루 + GMS(gpt-5-nano)\n서울 355개관 / 2153권 데이터", Inches(8.7), Inches(5.0), Inches(4.3), Inches(1.2), size=14, color=C_GRAY)

def slide_toc():
    s = prs.slides.add_slide(BLANK)
    bg_dark(s)
    accent_bar(s, C_POINT)
    add_text(s, "발표 순서", Inches(0.6), Inches(0.3), Inches(6), Inches(0.7), size=14, color=C_POINT, bold=True)
    add_text(s, "목차", Inches(0.6), Inches(0.9), Inches(8), Inches(0.9), size=38, bold=True)

    items = [
        ("01", "문제 정의 & 차별점", "독서하는 사람들이 겪는 불편과 '책잇다'의 솔루션"),
        ("02", "기술 스택 & 데이터", "Django · Vue · Leaflet 지도 · 2,153권 스냅샷"),
        ("03", "AI 추천 원리", "하이브리드 엔진 · 임베딩 유사도 · 환각 차단 로직"),
        ("04", "팀 구성 & 일정", "역할 분담 · 총 10일의 치열한 스프린트"),
        ("05", "시연 시나리오", "페르소나 기반 라이브 데모 (기획 의도 중심)"),
    ]
    for i, (num, title, sub) in enumerate(items):
        lft = Inches(0.5)
        top = Inches(2.0 + i * 1.0)
        card = add_rect(s, lft, top, Inches(12.33), Inches(0.85), fill=C_ACCENT, line=C_POINT)
        add_text(s, num, lft + Inches(0.2), top + Inches(0.15), Inches(0.8), Inches(0.5), size=24, bold=True, color=C_GOLD)
        add_text(s, title, lft + Inches(1.0), top + Inches(0.15), Inches(3.5), Inches(0.55), size=20, bold=True)
        add_text(s, sub, lft + Inches(4.5), top + Inches(0.22), Inches(7.5), Inches(0.55), size=15, color=C_GRAY)

def slide_problem_diff():
    s = prs.slides.add_slide(BLANK)
    bg_dark(s)
    accent_bar(s, C_POINT)
    add_text(s, "01  문제 정의 & 차별점", Inches(0.6), Inches(0.25), Inches(8), Inches(0.5), size=13, color=C_POINT, bold=True)
    add_text(s, "추천은 넘치지만, '지금 내 도서관'에 그 책이 없다", Inches(0.6), Inches(0.7), Inches(12), Inches(1.3), size=32, bold=True)

    # Pain points
    pains = [
        ("😫 완독 직후", "다음 뭐 읽지? 흐름이 끊긴다."),
        ("🏛️ 도서관에서", "무엇을 빌릴지 막막하고 재고가 없다."),
    ]
    for i, (label, desc) in enumerate(pains):
        card = add_rect(s, Inches(0.5 + i*6.3), Inches(1.8), Inches(6.0), Inches(1.5), fill=C_MID, line=C_RED)
        add_text(s, label, Inches(0.7 + i*6.3), Inches(2.0), Inches(5.6), Inches(0.5), size=18, bold=True, color=C_GOLD)
        add_text(s, desc, Inches(0.7 + i*6.3), Inches(2.5), Inches(5.6), Inches(0.5), size=15, color=C_LIGHT)

    # Differentiation
    add_text(s, "💡 책잇다만의 차별점", Inches(0.6), Inches(3.7), Inches(8), Inches(0.5), size=18, color=C_POINT, bold=True)
    cols = [
        ("기존 서점/AI", ["✅ 책 추천", "✅ 자연어 이유", "❌ 내 도서관 모름", "❌ 실시간 가용성 모름", "❌ 없는 책(환각) 추천 위험"], False),
        ("📚 책잇다", ["✅ 다중 seed 맞춤 AI 추천", "✅ 지금 빌릴 수 있는 책만", "✅ Leaflet 도서관 지도 지원", "✅ 추천 환각 완벽 차단"], True),
    ]
    for i, (name, points, highlight) in enumerate(cols):
        l = Inches(0.5 + i * 6.3)
        t = Inches(4.2)
        fill, line = (C_ACCENT, C_GOLD) if highlight else (C_MID, C_GRAY)
        card = add_rect(s, l, t, Inches(6.0), Inches(2.8), fill=fill, line=line)
        add_text(s, name, l+Inches(0.2), t+Inches(0.15), Inches(5.6), Inches(0.5), size=18 if highlight else 16, bold=True, color=C_GOLD if highlight else C_WHITE)
        for j, pt in enumerate(points):
            add_text(s, pt, l+Inches(0.2), t+Inches(0.75+j*0.4), Inches(5.6), Inches(0.4), size=14, color=C_WHITE if highlight or pt.startswith("✅") else C_GRAY)

def slide_tech():
    s = prs.slides.add_slide(BLANK)
    bg_dark(s)
    accent_bar(s, C_POINT)
    add_text(s, "02  기술 스택 & 데이터", Inches(0.6), Inches(0.25), Inches(8), Inches(0.5), size=13, color=C_POINT, bold=True)
    add_text(s, "검증된 스택과 2,153권의 풍부한 스냅샷 데이터", Inches(0.6), Inches(0.7), Inches(12), Inches(0.85), size=34, bold=True)

    stacks = [
        ("🖥️ Backend", "Django · DRF · SQLite", "dj-rest-auth 인증\n온디맨드 가용성(bookExist lazy)", C_ACCENT),
        ("🎨 Frontend", "Vue 3 · Pinia · Vite", "반응형 UI (미니멀/타이포 리디자인)\n읽음 슬라이드 토글", C_ACCENT),
        ("🤖 AI & 맵", "GMS · Leaflet · 임베딩", "text-embedding-3-small (내용유사도)\nOSM + Leaflet 도서관 지도 일원화", C_MID),
        ("📊 데이터", "정보나루 · 알라딘 API", "2153권 · 355관 · 2125개 신호\ndumpdata로 fixtures 재동결 완비", C_MID),
    ]
    for i, (label, name, desc, fill) in enumerate(stacks):
        l = Inches(0.5 + (i%2) * 6.3)
        t = Inches(1.9 + (i//2) * 2.5)
        card = add_rect(s, l, t, Inches(6.0), Inches(2.2), fill=fill, line=C_POINT)
        add_text(s, label, l+Inches(0.2), t+Inches(0.15), Inches(5.6), Inches(0.5), size=16, color=C_POINT, bold=True)
        add_text(s, name,  l+Inches(0.2), t+Inches(0.7), Inches(5.6), Inches(0.55), size=22, bold=True, color=C_WHITE)
        add_text(s, desc,  l+Inches(0.2), t+Inches(1.3), Inches(5.6), Inches(0.85), size=14, color=C_LIGHT)
        
    add_rect(s, Inches(0.5), Inches(6.8), Inches(12.33), Inches(0.5), fill=C_MID)
    add_text(s, "💡 외부 의존도 최소화: 데이터는 fixtures로 즉시 복원 가능. (데모 시 지도/가용성만 외부 통신)", Inches(0.7), Inches(6.85), Inches(12), Inches(0.4), size=13, color=C_GOLD)

def slide_arch():
    s = prs.slides.add_slide(BLANK)
    bg_dark(s)
    accent_bar(s, C_GOLD)
    add_text(s, "03  AI 추천 원리", Inches(0.6), Inches(0.25), Inches(8), Inches(0.5), size=13, color=C_GOLD, bold=True)
    add_text(s, "하이브리드 엔진: 데이터 필터링 + LLM 선별", Inches(0.6), Inches(0.7), Inches(12), Inches(0.85), size=32, bold=True)

    steps = [
        ("① 취향 시드", "온보딩 표지픽(다중 seed)\n+ 임베딩 발견", C_POINT),
        ("② 후보 생성", "함께대출/마니아/베스트셀러\n→ 대출가능한 책 필터", C_ACCENT),
        ("③ LLM 선별", "GMS gpt-5-nano\n(후보 풀 내에서만 선별)", C_MID),
        ("④ 후처리", "추천 결과 검증 및\n실제 데이터 기반 이유 생성", C_GOLD),
    ]
    for i, (num, desc, fill) in enumerate(steps):
        l = Inches(0.4 + i * 3.17)
        t = Inches(1.8)
        box = add_rect(s, l, t, Inches(2.8), Inches(1.6), fill=fill)
        num_color = C_DARK if fill == C_GOLD else C_WHITE
        add_text(s, num,  l+Inches(0.1), t+Inches(0.1), Inches(2.6), Inches(0.5), size=16, bold=True, color=num_color, align=PP_ALIGN.CENTER)
        add_text(s, desc, l+Inches(0.1), t+Inches(0.6), Inches(2.6), Inches(0.8), size=13, color=num_color, align=PP_ALIGN.CENTER)
        if i < 3: add_text(s, "→", Inches(3.25 + i*3.17), Inches(2.2), Inches(0.5), Inches(0.5), size=24, bold=True, color=C_POINT, align=PP_ALIGN.CENTER)

    rules = [
        ("🔒 LLM 환각 완벽 차단", "AI가 임의의 책을 지어내는 현상(환각)을 막기 위해, DB에 있는 후보 리스트를 제공하고\n해당 번호 안에서만 고르도록 강제하는 프롬프팅 적용 (후처리로 ISBN 2차 검증)"),
        ("⚡ 추론 속도 및 안정성", "temperature를 제거하고 reasoning_effort='low'로 설정하여 응답 속도 최적화\n상투적인 멘트를 차단하고 '좋아요'한 책 제목을 직접 인용하는 구체적 근거 생성"),
        ("💡 텍스트 임베딩 탐색", "text-embedding-3-small 내용 기반 임베딩 코사인 유사도로 '취향 발견' 기능 제공"),
    ]
    for i, (title, desc) in enumerate(rules):
        t = Inches(3.7 + i * 1.1)
        card = add_rect(s, Inches(0.5), t, Inches(12.33), Inches(0.9), fill=C_ACCENT, line=C_POINT)
        add_text(s, title, Inches(0.7), t+Inches(0.15), Inches(3.5), Inches(0.6), size=16, bold=True, color=C_GOLD)
        add_text(s, desc, Inches(3.8), t+Inches(0.15), Inches(8.8), Inches(0.6), size=13, color=C_LIGHT)

def slide_team():
    s = prs.slides.add_slide(BLANK)
    bg_dark(s)
    accent_bar(s, C_GOLD)
    add_text(s, "04  팀 구성 & 역할 분담", Inches(0.6), Inches(0.25), Inches(8), Inches(0.5), size=13, color=C_GOLD, bold=True)
    add_text(s, "2인 팀 — 기능별 명확한 분담, 추천/통합은 공동 작업", Inches(0.6), Inches(0.7), Inches(12), Inches(0.85), size=30, bold=True)

    card_a = add_rect(s, Inches(0.5), Inches(1.75), Inches(5.9), Inches(5.3), fill=C_ACCENT, line=C_GOLD)
    add_text(s, "👩‍💻  팀원 A", Inches(0.7), Inches(1.85), Inches(5.5), Inches(0.55), size=20, bold=True)
    fe_badge = add_rect(s, Inches(0.7), Inches(2.38), Inches(2.2), Inches(0.38), fill=C_GOLD)
    add_text(s, "Frontend 담당", Inches(0.7), Inches(2.4), Inches(2.2), Inches(0.35), size=13, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    tasks_a = [
        "Vue 3 + Vite 반응형 웹 UI 및 라우팅",
        "타이포그래피 및 미니멀 UI 리디자인",
        "온보딩 '표지픽' 취향 입력 컴포넌트",
        "임베딩 기반 홈 '취향 발견' 및 내 서재 통계",
        "책 상세 Leaflet 기반 도서관 지도 일원화",
        "커뮤니티(게시판) 게시글/댓글/좋아요 UI",
    ]
    for i, task in enumerate(tasks_a):
        add_text(s, f"• {task}", Inches(0.7), Inches(2.9 + i * 0.42), Inches(5.5), Inches(0.45), size=13, color=C_LIGHT)

    card_b = add_rect(s, Inches(6.8), Inches(1.75), Inches(5.9), Inches(5.3), fill=C_ACCENT, line=C_POINT)
    add_text(s, "👨‍💻  팀원 B", Inches(7.0), Inches(1.85), Inches(5.5), Inches(0.55), size=20, bold=True)
    be_badge = add_rect(s, Inches(7.0), Inches(2.38), Inches(2.2), Inches(0.38), fill=C_POINT)
    add_text(s, "Backend 담당", Inches(7.0), Inches(2.4), Inches(2.2), Inches(0.35), size=13, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    tasks_b = [
        "Django DB 모델링 및 dj-rest-auth 인증",
        "다중 seed 기반 후보 도서 생성(build_candidates)",
        "GMS LLM 프롬프트 설계 및 환각 차단 후처리",
        "온디맨드 가용성(bookExist lazy) TTL 캐시 처리",
        "커뮤니티 API 구현 (Article/Comment CRUD)",
        "서울 355개관 등 데이터 Fixtures 동결 (더미 포함)",
    ]
    for i, task in enumerate(tasks_b):
        add_text(s, f"• {task}", Inches(7.0), Inches(2.9 + i * 0.42), Inches(5.5), Inches(0.45), size=13, color=C_LIGHT)

    add_rect(s, Inches(0.5), Inches(7.15), Inches(12.33), Inches(0.22), fill=C_MID)
    add_text(s, "공동: 추천 엔진 아키텍처 설계 · 하이브리드 로직 통합 · 데모 시연 준비", Inches(0.7), Inches(7.17), Inches(12), Inches(0.2), size=12, color=C_GOLD)

def slide_schedule():
    s = prs.slides.add_slide(BLANK)
    bg_dark(s)
    accent_bar(s, C_POINT)
    add_text(s, "04  개발 일정 (총 10일)", Inches(0.6), Inches(0.25), Inches(8), Inches(0.5), size=13, color=C_POINT, bold=True)
    add_text(s, "10일 스프린트 — 기획부터 시연까지", Inches(0.6), Inches(0.7), Inches(12), Inches(0.75), size=30, bold=True)

    TOTAL_DAYS = 10
    LABEL_W = Inches(2.5)
    TRACK_L = Inches(2.7)
    TRACK_W = Inches(10.3)
    TRACK_H = Inches(0.5)
    ROW_GAP  = Inches(0.75)
    START_T  = Inches(1.8)
    
    for d in range(TOTAL_DAYS):
        add_text(s, f"Day {d+1}", TRACK_L + (TRACK_W / TOTAL_DAYS) * d, Inches(1.5), TRACK_W / TOTAL_DAYS, Inches(0.25), size=11, color=C_GRAY, align=PP_ALIGN.CENTER)

    rows = [
        ("Phase 1~2: 토대/수집",       0, 2,  C_MID),
        ("Phase 3: 기본 API/UI",       2, 3,  C_POINT),
        ("Phase 4~5: LLM 추천/서재",   3, 4,  C_POINT),
        ("Phase 6: 다중 seed/환각차단",4, 3,  C_GOLD),
        ("지도일원화/임베딩/리디자인", 6, 3,  C_GOLD),
        ("Fixtures 동결 및 발표준비",  8, 2,  C_RED),
    ]

    for ri, (label, start, dur, color) in enumerate(rows):
        t = START_T + ROW_GAP * ri
        add_rect(s, TRACK_L, t, TRACK_W, TRACK_H, fill=C_DARK, line=C_ACCENT)
        add_text(s, label, Inches(0.1), t+Inches(0.05), LABEL_W, TRACK_H, size=13, color=C_WHITE, align=PP_ALIGN.RIGHT)
        bar_l = TRACK_L + (TRACK_W / TOTAL_DAYS) * start
        bar_w = (TRACK_W / TOTAL_DAYS) * dur
        add_rect(s, bar_l, t + Emu(50000), bar_w, TRACK_H - Emu(100000), fill=color)

def slide_persona():
    s = prs.slides.add_slide(BLANK)
    bg_dark(s)
    accent_bar(s, C_POINT)
    add_text(s, "05  시연 시나리오", Inches(0.6), Inches(0.25), Inches(8), Inches(0.5), size=13, color=C_POINT, bold=True)
    add_text(s, "스토리텔링 기반 라이브 데모 (기획 의도 중심)", Inches(0.6), Inches(0.7), Inches(12), Inches(0.85), size=36, bold=True)

    pcard = add_rect(s, Inches(0.5), Inches(1.8), Inches(4.5), Inches(5.2), fill=C_ACCENT, line=C_GOLD)
    add_text(s, "👤  페르소나", Inches(0.7), Inches(1.95), Inches(4.1), Inches(0.5), size=14, bold=True, color=C_GOLD)
    add_text(s, "김싸피 씨\n(30대, 직장인)", Inches(0.7), Inches(2.45), Inches(4.1), Inches(0.9), size=24, bold=True)
    add_text(s, "• 퇴근 후 방금 책을 완독함.\n\n• '다음에 뭐 읽지?' 고민됨.\n\n• 마포 도서관을 종종 방문하지만, 막상 가면 빌릴 책이 없어 헛걸음함.\n\n• 긴 타이핑이나 뻔한 설명서 대신, 당장 빌릴 수 있는 해답을 원함.", Inches(0.7), Inches(3.6), Inches(4.1), Inches(3.0), size=14, color=C_LIGHT)

    scenario = add_rect(s, Inches(5.2), Inches(1.8), Inches(7.7), Inches(5.2), fill=C_MID, line=C_POINT)
    add_text(s, "📋  시연 흐름 (불필요한 과정 과감히 생략)", Inches(5.4), Inches(1.95), Inches(7.3), Inches(0.5), size=14, bold=True, color=C_POINT)
    flows = [
        ("Step 1", "취향 입력 (온보딩)",  "회원가입/로그인 생략. '표지픽'으로 텍스트 입력 없이 취향 셋업"),
        ("Step 2", "맞춤형 AI 추천",     "추천된 도서와 함께, '왜 이 책인지' 기획 의도와 직접적인 근거 제시"),
        ("Step 3", "취향 발견 (임베딩)",  "text-embedding-3-small 기반 '이런 책도 잇다' 확장 탐색"),
        ("Step 4", "실시간 도서관 매핑",  "책 상세 뷰의 Leaflet 지도에서 내 주변 도서관 실시간 가용성 확인"),
        ("Step 5", "커뮤니티 소통",       "가용성을 확인한 책을 완독 후, 게시판에서 서평과 댓글로 타인과 공유"),
        ("Step 6", "풍성한 내 서재",      "미리 채워둔 더미 데이터로 통계 및 좋아요 목록 시연"),
    ]
    for i, (step, title, desc) in enumerate(flows):
        t_row = Inches(2.5 + i * 0.75)
        add_text(s, step,  Inches(5.4), t_row, Inches(0.9), Inches(0.7), size=12, bold=True, color=C_GOLD)
        add_text(s, title, Inches(6.3), t_row, Inches(2.0), Inches(0.7), size=14, bold=True)
        add_text(s, desc,  Inches(8.3), t_row+Inches(0.03), Inches(4.4), Inches(0.7), size=12, color=C_GRAY)

def slide_closing():
    s = prs.slides.add_slide(BLANK)
    bg_dark(s)
    add_rect(s, 0, 0, W, H, fill=C_MID)
    accent_bar(s, C_GOLD, Inches(0.1))
    add_rect(s, 0, H-Inches(0.1), W, Inches(0.1), fill=C_GOLD)
    add_text(s, "책잇다", Inches(1), Inches(1.3), Inches(11.33), Inches(1.4), size=64, bold=True, align=PP_ALIGN.CENTER)

    m1 = add_rect(s, Inches(3.0), Inches(2.85), Inches(3.3), Inches(0.42), fill=C_ACCENT, line=C_GOLD)
    add_text(s, "책 · 도서관 · 사람을 잇다", Inches(3.0), Inches(2.88), Inches(3.3), Inches(0.38), size=14, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)
    m2 = add_rect(s, Inches(6.6), Inches(2.85), Inches(2.5), Inches(0.42), fill=C_ACCENT, line=C_POINT)
    add_text(s, "도서관에 책이 있다", Inches(6.6), Inches(2.88), Inches(2.5), Inches(0.38), size=14, bold=True, color=C_POINT, align=PP_ALIGN.CENTER)

    add_text(s, "완독 직후, 도서관 방문 전 — 딱 그 순간\n지금 내 도서관에서 빌릴 수 있는 다음 한 권", Inches(1), Inches(3.45), Inches(11.33), Inches(1.1), size=19, color=C_LIGHT, align=PP_ALIGN.CENTER)
    summary = ["데이터(정보나루 실제 대출)로 후보를 만들고", "AI(GMS)가 후보 안에서만 이유와 함께 선별하며", "지금 내 도서관에서 빌릴 수 있는 책만 보여준다"]
    for i, line in enumerate(summary):
        add_text(s, f"  {i+1}.  {line}", Inches(2.5), Inches(4.75+i*0.6), Inches(8.5), Inches(0.55), size=16, color=C_GOLD if i == 2 else C_LIGHT)
    add_text(s, "감사합니다  🙏     Q & A", Inches(1), Inches(6.5), Inches(11.33), Inches(0.85), size=28, bold=True, color=C_POINT, align=PP_ALIGN.CENTER)

slide_cover()
slide_toc()
slide_problem_diff()
slide_tech()
slide_arch()
slide_team()
slide_schedule()
slide_persona()
slide_closing()

OUTPUT = "책잇다_발표.pptx"
prs.save(OUTPUT)
print(f"[완료] 저장: {OUTPUT}  ({prs.slides.__len__()}장)")
